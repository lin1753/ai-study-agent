# parsing.py
import os
import logging
from pypdf import PdfReader
from pptx import Presentation

# Mute noisy pdfminer warnings (like FontBBox issues)
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logger = logging.getLogger(__name__)

import re
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

def clean_text(text: str) -> str:
    """
    Clean text extracted from PDF/PPT.
    1. Merge broken lines (e.g. "Net\nwork" -> "Network").
    2. Fix vertical text (e.g. "N\ne\nt" -> "Net").
    3. Remove excessive whitespace.
    """
    if not text:
        return ""
    
    encoded_text = text.encode("utf-8", "ignore").decode("utf-8")
    
    # Debugging: Log raw length
    raw_len = len(encoded_text)
    
    # Simplified Cleaning:
    # Replace newlines with spaces to keep word boundaries
    text = encoded_text.strip().replace('\n', ' ')
    
    # Remove excessive spaces
    text = re.sub(r' +', ' ', text)
    
    clean_len = len(text)
    # print(f"[DEBUG] Cleaned: {raw_len} -> {clean_len} chars")
    
    return text

def parse_pdf(filepath: str, user_config: dict = None) -> list[str]:
    """
    Parse PDF and return a list of cleaned text blocks.
    Strategy: Use pdfplumber primarily (better layout handling). Fallback to pypdf.
    """
    blocks = []
    print(f"[DEBUG] Parsing PDF: {filepath}")
    
    # Method 1: pdfplumber (Primary)
    if pdfplumber:
        print("[INFO] Using pdfplumber as primary parser...")
        try:
            with pdfplumber.open(filepath) as pdf:
                print(f"[DEBUG] Total Pages (pdfplumber): {len(pdf.pages)}")
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    # print(f"[DEBUG] pdfplumber Page {i} Raw: {text[:50]!r}")
                    cleaned = clean_text(text)
                    if cleaned:
                        # print(f"[DEBUG] pdfplumber Page {i} Cleaned (First 100): {cleaned[:100]!r}")
                        blocks.append(cleaned)
        except Exception as e:
            print(f"[ERROR] pdfplumber parsing failed: {e}")
    else:
        print("[WARN] pdfplumber not installed. Skipping.")

    # Quality Check 1: pdfplumber
    total_len = sum(len(b) for b in blocks)
    if len(blocks) > 0 and total_len >= 500:
        print(f"[DEBUG] pdfplumber successful. Extracted {len(blocks)} blocks, total chars: {total_len}")
        return blocks
    else:
        print(f"[INFO] pdfplumber result insufficient (blocks={len(blocks)}, chars={total_len}). Continuing to fallback...")
        # Keep blocks, maybe pypdf will add more

    # Method 2: pypdf (Fallback)
    if not blocks:
        print("[INFO] Falling back to pypdf...")
        try:
            reader = PdfReader(filepath)
            print(f"[DEBUG] Total Pages (pypdf): {len(reader.pages)}")
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                cleaned = clean_text(text)
                if cleaned:
                    blocks.append(cleaned)
        except Exception as e:
            print(f"[ERROR] pypdf parsing failed: {e}")

    # Quality Check 2: pypdf & Final OCR Decision
    total_len = sum(len(b) for b in blocks)
    logger.info(f"[DEBUG] Total text length from extractors: {total_len}")
    
    # If the document yields less than 1000 characters (e.g. image-heavy slides), fallback to OCR
    if total_len < 1000:
        logger.warning(f"[WARN] Text extraction yielded only {total_len} chars (< 1000). Attempting RapidOCR via PyMuPDF...")
        ocr_blocks = parse_pdf_ocr(filepath) 
        if ocr_blocks:
            blocks.extend(ocr_blocks)
        total_len = sum(len(b) for b in blocks)
        logger.info(f"[DEBUG] Total text length after OCR fallback: {total_len}")

    logger.info(f"[DEBUG] Final Extracted {len(blocks)} non-empty blocks.")
    return blocks

def parse_pdf_ocr(filepath: str) -> list[str]:
    """
    Fallback method: Convert PDF pages to images using PyMuPDF and run pure-Python OCR (RapidOCR) 
    to extract text from heavily image-based documents. Very fast on CPU.
    """
    blocks = []
    
    try:
        import fitz  # PyMuPDF
        import numpy as np
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as e:
        logger.error(f"[OCR] Required packages not found. Please run 'pip install pymupdf rapidocr_onnxruntime'. Error: {e}")
        return []

    try:
        logger.info("[OCR] Initializing RapidOCR ONNX model...")
        # Initialize OCR engine
        engine = RapidOCR()
        
        doc = fitz.open(filepath)
        logger.info(f"[OCR] Total {len(doc)} pages detected. Starting OCR extraction...")
        
        # We limit the max pages just in case to prevent infinite hangs, though RapidOCR is very fast.
        max_pages = min(len(doc), 50) 
        
        for i in range(max_pages):
            page = doc[i]
            logger.debug(f"[OCR] Processing page {i+1}/{max_pages}...")
            
            # Render page to zoom 2.0 (around 144 DPI) which is a sweet spot for RapidOCR
            zoom_matrix = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=zoom_matrix, alpha=False)
            
            # Convert PyMuPDF unformatted raw bytes to NumPy array (H, W, Channels)
            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
            
            # Run OCR prediction on the numpy image
            # Result format: List of [ [ [x, y], ...], "Extracted Text", Confidence ]
            ocr_result, elapse = engine(img_array)
            
            if ocr_result:
                # Combine all text blocks detected on this page
                page_text = "\n".join([item[1] for item in ocr_result])
                cleaned = clean_text(page_text)
                if cleaned:
                    blocks.append(cleaned)
                    
        logger.info(f"[OCR] Completed processing {max_pages} pages.")

    except Exception as e:
        logger.error(f"[OCR] Fatal error during RapidOCR execution: {e}", exc_info=True)
        
    return blocks

def parse_ppt(filepath: str) -> list[str]:
    """
    Parse PPTX and return a list of text blocks (slides).
    """
    blocks = []
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            
            full_text = "\n".join(text_runs).strip()
            if full_text:
                blocks.append(full_text)
    except Exception as e:
        print(f"Error parsing PPT {filepath}: {e}")
        return []
    return blocks

def parse_file(filepath: str, file_type: str, user_config: dict = None) -> list[str]:
    if file_type.lower() == 'pdf':
        return parse_pdf(filepath, user_config=user_config)
    elif file_type.lower() in ['ppt', 'pptx']:
        return parse_ppt(filepath)
    else:
        return []
